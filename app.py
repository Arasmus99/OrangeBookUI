import re
import pandas as pd
from pptx import Presentation
from datetime import date, timedelta
from dateutil.parser import parse
import streamlit as st
from io import BytesIO

# === Regex Patterns ===
PATTERNS = {
    "docket_number": re.compile(
        r"\b\d{4}-[A-Z]{2,}-\d{5}-\d{2}\b"           # e.g. 2018-LOW-68327-13
        r"|\b\d{5}-\d{2}\b"                           # e.g. 68327-13
        r"|\b\d{5}-\d{4}-\d{2}[A-Z]{2,4}\b"           # e.g. 01330-0004-00US
        r"|\b\d{4}\.\d{3}-?[A-Z]{2,}\d*\b"          # e.g. 0509.003US, 0509.003-US, 0509.0003-US8
        r"|\b\d{4}-\d{4}-[A-Z]{3}\b"                  # e.g. 0509-0001-PCT
    ),
    "application_number": re.compile(r"\b\d{2}/\d{3}[,]?\d{3}\s+[A-Z]{2}\b"),
    "alt_application_number": re.compile(
        r"\b[Pp]\d{11}\s+[A-Z]{2}-\w{1,4}\b"
        r"|\b\d{5,12}(?:[.,]\d+)?\s+[A-Z]{2,3}\b"
    ),
    "pct_number": re.compile(r"PCT/[A-Z]{2}\d{4}/\d{6}"),
    "wipo_number": re.compile(r"\bWO\d{4}/\d{6}\b"),
    "date": re.compile(r"\b(\d{1,2}[/-]\d{1,2}[/-]\d{2,4})\b")
}

SKIP_PHRASES = ["PENDING", "ABANDONED", "WITHDRAWN", "GRANTED", "ISSUED", "STRUCTURE"]

# === Recursive text extraction for GroupShapes ===
def extract_texts_from_shape_recursive(shape):
    texts = []
    if shape.shape_type == 6:  # GroupShape
        for shp in shape.shapes:
            texts.extend(extract_texts_from_shape_recursive(shp))
    else:
        text = extract_text_from_shape(shape)
        if text:
            texts.append(text)
    return texts

def extract_text_from_shape(shape):
    if shape.has_text_frame:
        return shape.text.strip()
    return ""

def should_include(text):
    upper_text = text.upper()
    return not any(phrase in upper_text for phrase in SKIP_PHRASES)

def get_earliest_due_date(dates_str):
    if not isinstance(dates_str, str):
        return pd.NaT
    try:
        dates = [parse(d.strip(), dayfirst=False, fuzzy=True) for d in dates_str.split(";") if d.strip()]
        return min(dates) if dates else pd.NaT
    except:
        return pd.NaT

def extract_entries_from_textbox(text, months_back=0):
    entries = []
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    cutoff_date = date.today() - timedelta(days=30 * months_back)

    entry = {
        "docket_number": None,
        "application_number": None,
        "pct_number": None,
        "wipo_number": None,
        "due_dates": [],
        "raw_text": "\n".join(lines)
    }

    for line in lines:
        clean_line = line.replace(" /,", "/").replace("/", "/").replace(",,", ",").replace(" /", "/")
        clean_line = re.sub(r"[^0-9A-Za-z/,\.\s-]", "", clean_line)  # Remove odd characters
        clean_line = clean_line.replace(",", "")  # Remove commas inside numbers like 201,375

        if not entry["docket_number"] and PATTERNS["docket_number"].search(clean_line):
            entry["docket_number"] = PATTERNS["docket_number"].search(clean_line).group(0)

        if not entry["pct_number"] and PATTERNS["pct_number"].search(clean_line):
            entry["pct_number"] = PATTERNS["pct_number"].search(clean_line).group(0)

        if not entry["application_number"] and PATTERNS["application_number"].search(clean_line):
            entry["application_number"] = PATTERNS["application_number"].search(clean_line).group(0)
        elif not entry["application_number"] and PATTERNS["alt_application_number"].search(clean_line):
            entry["application_number"] = PATTERNS["alt_application_number"].search(clean_line).group(0)

        if not entry["wipo_number"] and PATTERNS["wipo_number"].search(clean_line):
            entry["wipo_number"] = PATTERNS["wipo_number"].search(clean_line).group(0)

        for match in PATTERNS["date"].findall(clean_line):
            try:
                parsed = parse(match, dayfirst=False, fuzzy=True)
                if parsed.date() >= cutoff_date:
                    entry["due_dates"].append(parsed.strftime("%m/%d/%Y"))
            except:
                continue

    if not (entry["docket_number"] or entry["application_number"] or entry["pct_number"] or entry["wipo_number"]):
        return []

    if entry["due_dates"]:
        entries.append(entry)

    return entries

def extract_from_pptx(upload, months_back):
    prs = Presentation(upload)
    results = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape_num, shape in enumerate(slide.shapes, start=1):
            texts = extract_texts_from_shape_recursive(shape)
            for text in texts:
                if not should_include(text):
                    continue
                entries = extract_entries_from_textbox(text, months_back)
                for entry in entries:
                    results.append({
                        "Slide": slide_num,
                        "Textbox Content": entry["raw_text"],
                        "Docket Number": entry["docket_number"],
                        "Application Number": entry["application_number"],
                        "PCT Number": entry["pct_number"],
                        "WIPO Number": entry["wipo_number"],
                        "Due Dates": "; ".join(entry["due_dates"])
                    })

    if not results:
        return pd.DataFrame(columns=["Slide", "Textbox Content", "Docket Number", "Application Number", "PCT Number", "WIPO Number", "Due Dates"])

    df = pd.DataFrame(results)
    df["Earliest Due Date"] = df["Due Dates"].apply(get_earliest_due_date)
    df = df.sort_values(by="Earliest Due Date", ascending=True).drop(columns=["Earliest Due Date"])
    return df

# === Streamlit UI ===
st.title("\U0001F4CA PowerPoint Patent Extractor")
ppt_files = st.file_uploader("Upload one or more PowerPoint (.pptx) files", type="pptx", accept_multiple_files=True)
months_back = st.slider("Include due dates up to this many months in the past:", 0, 24, 0)

if ppt_files:
    all_dfs = []
    for ppt_file in ppt_files:
        df = extract_from_pptx(ppt_file, months_back)
        if df.empty:
            st.warning(f"⚠️ No extractable data found in {ppt_file.name}.")
            continue
        df["Filename"] = ppt_file.name
        all_dfs.append(df)

    if all_dfs:
        final_df = pd.concat(all_dfs, ignore_index=True)
        st.success(f"✅ Extracted {len(final_df)} entries from {len(all_dfs)} file(s).")
        st.dataframe(final_df, use_container_width=True)

        output = BytesIO()
        final_df.to_excel(output, index=False)
        output.seek(0)
        st.download_button("\U0001F4E5 Download Excel", output, file_name="combined_extracted_data.xlsx")
